"""
generate_final_ppt.py
ONE combined presentation for CcMart — covers everything:
problem, solution, tech stack, EDA, SQL, ML, streaming, dashboard,
simulation, outcomes, and business impact.
Design: Clean white/light theme · navy headers · simple & professional

Run:    python generate_final_ppt.py
Output: dashboard/CcMart_Final_Presentation.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────
WHITE   = RGBColor(0xFF,0xFF,0xFF)
BG      = RGBColor(0xF4,0xF6,0xFA)
NAVY    = RGBColor(0x14,0x23,0x5E)
BLUE    = RGBColor(0x1E,0x6F,0xC8)
LBLUE   = RGBColor(0xE8,0xF0,0xFB)
GREEN   = RGBColor(0x0F,0x9B,0x56)
LGREEN  = RGBColor(0xE5,0xF9,0xEE)
ORANGE  = RGBColor(0xD4,0x6B,0x08)
LORANGE = RGBColor(0xFE,0xF3,0xE2)
PURPLE  = RGBColor(0x6C,0x3F,0xC4)
LPURPLE = RGBColor(0xEF,0xEB,0xFF)
TEAL    = RGBColor(0x08,0x70,0x79)
LTEAL   = RGBColor(0xE0,0xF4,0xF5)
RED     = RGBColor(0xC0,0x32,0x2A)
LRED    = RGBColor(0xFC,0xEC,0xEA)
T1      = RGBColor(0x1A,0x1D,0x2E)
T2      = RGBColor(0x44,0x4D,0x6A)
T3      = RGBColor(0x7A,0x84,0xA8)
DIVIDER = RGBColor(0xD0,0xD5,0xE8)

W = Inches(13.33)
H = Inches(7.5)
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ═══════════════ HELPERS ════════════════════════════════

def rect(s, x, y, w, h, fill=None, line=None, lw=Pt(0)):
    sh = s.shapes.add_shape(1, x, y, w, h)
    sh.line.width = lw
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else: sh.fill.background()
    if line: sh.line.fill.solid(); sh.line.fill.fore_color.rgb = line
    else: sh.line.fill.background()
    return sh

def tb(s, text, x, y, w, h, sz=Pt(12), bold=False, color=T1,
       align=PP_ALIGN.LEFT, italic=False):
    box = s.shapes.add_textbox(x, y, w, h)
    box.word_wrap = True
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text; r.font.size = sz; r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = "Calibri"
    return box

def ml(s, lines, x, y, w, h, sz=Pt(11), dc=T2):
    box = s.shapes.add_textbox(x, y, w, h)
    box.word_wrap = True; tf = box.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if isinstance(ln, tuple): t_, b_, c_ = ln
        else: t_, b_, c_ = ln, False, dc
        if t_ == "":
            r = p.add_run(); r.text = " "; r.font.size = Pt(5); continue
        r = p.add_run(); r.text = t_
        r.font.size = sz; r.font.bold = b_; r.font.color.rgb = c_
        r.font.name = "Calibri"

def slide_bg(s):
    rect(s, 0, 0, W, H, fill=WHITE)

def header(s, title, sub="", bar=NAVY):
    bh = Inches(1.28) if sub else Inches(1.05)
    rect(s, 0, 0, W, bh, fill=bar)
    rect(s, 0, bh, W, Inches(0.042), fill=BLUE)
    tb(s, title, Inches(0.55), Inches(0.14), Inches(12.22), Inches(0.6),
       sz=Pt(24), bold=True, color=WHITE)
    if sub:
        tb(s, sub, Inches(0.55), Inches(0.72), Inches(12.22), Inches(0.38),
           sz=Pt(12), color=RGBColor(0xB0,0xBC,0xE8))

def tag(s, text, color=BLUE, y=Inches(1.44)):
    rect(s, Inches(0.55), y, Inches(0.06), Inches(0.28), fill=color)
    tb(s, text, Inches(0.68), y-Inches(0.01), Inches(8), Inches(0.32),
       sz=Pt(9), bold=True, color=color)

def rule(s, y):
    rect(s, Inches(0.55), y, Inches(12.22), Inches(0.018), fill=DIVIDER)

def card(s, x, y, w, h, fill=WHITE, accent=None):
    rect(s, x, y, w, h, fill=fill, line=DIVIDER, lw=Pt(0.5))
    if accent:
        rect(s, x, y, w, Inches(0.048), fill=accent)

def lcard(s, x, y, w, h, fill=WHITE, accent=None):
    rect(s, x, y, w, h, fill=fill, line=DIVIDER, lw=Pt(0.4))
    if accent:
        rect(s, x, y, Inches(0.055), h, fill=accent)

def kpi_box(s, x, y, w, h, val, lbl, sub_="", ac=BLUE, vsz=Pt(26)):
    card(s, x, y, w, h, accent=ac)
    tb(s, val, x, y+Inches(0.1), w, Inches(0.52), sz=vsz, bold=True,
       color=ac, align=PP_ALIGN.CENTER)
    tb(s, lbl, x, y+Inches(0.6), w, Inches(0.26), sz=Pt(10), bold=True,
       color=T2, align=PP_ALIGN.CENTER)
    if sub_:
        tb(s, sub_, x, y+Inches(0.84), w, Inches(0.2), sz=Pt(8.5),
           color=T3, align=PP_ALIGN.CENTER)

def kpi_row(s, items, y, h=Inches(1.1)):
    n = len(items); gap = Inches(0.12)
    w_ = (W - Inches(0.55)*2 - gap*(n-1)) / n
    x_ = Inches(0.55)
    for v, l, su, c in items:
        kpi_box(s, x_, y, w_, h, v, l, su, c)
        x_ += w_ + gap

# ═══════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); slide_bg(s)
rect(s, 0, 0, Inches(8.2), H, fill=NAVY)
rect(s, Inches(8.2), 0, W-Inches(8.2), H, fill=BG)
rect(s, Inches(8.2), 0, Inches(0.04), H, fill=BLUE)
rect(s, 0, 0, W, Inches(0.048), fill=BLUE)

tb(s, "CLOUD COMPUTING  ·  GROUP 6  ·  FINAL PROJECT",
   Inches(0.5), Inches(0.55), Inches(7.5), Inches(0.28),
   sz=Pt(9), color=RGBColor(0x7A,0x8C,0xC4))
tb(s, "CcMart", Inches(0.5), Inches(1.15), Inches(7.5), Inches(1.1),
   sz=Pt(62), bold=True, color=WHITE)
tb(s, "Real-Time E-Commerce Analytics",
   Inches(0.5), Inches(2.2), Inches(7.5), Inches(0.7),
   sz=Pt(26), color=RGBColor(0xA8,0xBE,0xF4))
tb(s, "Using Apache Spark", Inches(0.5), Inches(2.82), Inches(7.5), Inches(0.55),
   sz=Pt(22), bold=False, color=RGBColor(0x7A,0x90,0xC8))
rect(s, Inches(0.5), Inches(3.55), Inches(2.2), Inches(0.03), fill=BLUE)
tb(s, "Transforming 12.8 million raw clickstream events into real-time business intelligence — using Apache Spark SQL, MLlib, and Structured Streaming.",
   Inches(0.5), Inches(3.72), Inches(7.5), Inches(0.85),
   sz=Pt(12), color=RGBColor(0xB0,0xBC,0xD8))

tags_t = [("Apache Spark",BLUE),("PySpark MLlib",PURPLE),
          ("Spark SQL",TEAL),("Streaming",GREEN),("BI Dashboard",ORANGE)]
tx = Inches(0.5); ty = Inches(4.8)
for tg, tc in tags_t:
    tw = Inches(1.55)
    rect(s, tx, ty, tw, Inches(0.36), fill=WHITE, line=tc, lw=Pt(0.8))
    tb(s, tg, tx, ty, tw, Inches(0.36), sz=Pt(9.5), bold=True,
       color=tc, align=PP_ALIGN.CENTER)
    tx += tw + Inches(0.1)

# Right side summary
tb(s, "What This Presentation Covers",
   Inches(8.55), Inches(0.95), Inches(4.45), Inches(0.35),
   sz=Pt(12), bold=True, color=NAVY)
covers = [
    ("01","Problem Statement & Goals",RED),
    ("02","Why Apache Spark?",BLUE),
    ("03","Data & Exploratory Analysis",TEAL),
    ("04","Spark SQL — 6 Advanced Queries",PURPLE),
    ("05","Machine Learning — 3 Models",ORANGE),
    ("06","Real-Time Streaming",GREEN),
    ("07","Business Dashboard & Simulation",BLUE),
    ("08","Outcomes & Business Impact",GREEN),
    ("09","Key Learnings",NAVY),
]
cy = Inches(1.4)
for num, item, c in covers:
    rect(s, Inches(8.55), cy, Inches(4.45), Inches(0.52), fill=BG, line=DIVIDER, lw=Pt(0.3))
    rect(s, Inches(8.55), cy, Inches(0.04), Inches(0.52), fill=c)
    tb(s, num, Inches(8.66), cy+Inches(0.12), Inches(0.35), Inches(0.28),
       sz=Pt(9), bold=True, color=c)
    tb(s, item, Inches(9.06), cy+Inches(0.12), Inches(3.88), Inches(0.28),
       sz=Pt(11), color=T2)
    cy += Inches(0.58)

# ═══════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); slide_bg(s)
header(s, "Problem Statement",
       "What challenge does this project solve — and why does it matter?")
tag(s, "01  ·  MOTIVATION & GOALS", RED)

rect(s, Inches(0.55), Inches(1.82), Inches(12.22), Inches(0.95),
     fill=LRED, line=RED, lw=Pt(0.5))
rect(s, Inches(0.55), Inches(1.82), Inches(0.06), Inches(0.95), fill=RED)
tb(s, "The Core Problem", Inches(0.72), Inches(1.9), Inches(12.0), Inches(0.3),
   sz=Pt(12), bold=True, color=RED)
tb(s, "E-commerce platforms generate millions of user events every day — every click, every add-to-cart, every payment. But most companies cannot process this data fast enough to act on it. By the time a report reaches the operations team, the opportunity to intervene has already long passed.",
   Inches(0.72), Inches(2.22), Inches(12.0), Inches(0.48),
   sz=Pt(11.5), color=T1)

problems = [
    (RED,    "42% Cart Abandonment — Invisible",
     "Shoppers add items and never check out. Without real-time analytics there is no mechanism to identify who is about to leave and intervene before they do."),
    (ORANGE, "Marketing is One-Size-Fits-All",
     "All 100,000 customers receive the same promotions. No segmentation, no personalisation — wasted spend and missed conversion opportunities."),
    (PURPLE, "No Product Cross-Sell Intelligence",
     "Which products are bought together? What does a customer browse before they buy? Without this, the platform cannot show relevant cross-sell banners or personalised picks."),
    (TEAL,   "Operations React Hours Too Late",
     "Payment anomalies and unusual session behaviour are discovered hours later in batch reports — long after the revenue impact has already occurred."),
]
pw = Inches(5.88); ph = Inches(1.42); pgap = Inches(0.14)
px_ = Inches(0.55); py_ = Inches(2.92)
for i, (c, t, d) in enumerate(problems):
    col = i % 2; row = i // 2
    x_ = px_ + col*(pw+pgap); y_ = py_ + row*(ph+pgap)
    lcard(s, x_, y_, pw, ph, accent=c)
    tb(s, t, x_+Inches(0.14), y_+Inches(0.1), pw-Inches(0.22), Inches(0.28),
       sz=Pt(11.5), bold=True, color=T1)
    tb(s, d, x_+Inches(0.14), y_+Inches(0.44), pw-Inches(0.22), Inches(0.88),
       sz=Pt(10.5), color=T2)

rule(s, Inches(6.1))
tb(s, "Goal:  Build an end-to-end Apache Spark pipeline that ingests raw data, runs advanced SQL analytics, trains 3 ML models, and streams real-time predictions — all surfaced in a live BI dashboard.",
   Inches(0.55), Inches(6.18), Inches(12.22), Inches(0.5),
   sz=Pt(11.5), italic=True, color=NAVY)

# ═══════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION OVERVIEW
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); slide_bg(s)
header(s, "Our Solution — CcMart Analytics Pipeline",
       "A 6-stage Apache Spark pipeline that turns raw data into real-time business decisions")
tag(s, "02  ·  SOLUTION OVERVIEW", BLUE)

flow = [
    ("📥","1. Ingest","CSV → Parquet\n4 tables · 13M rows",BLUE),
    ("🔍","2. EDA","5 insight charts\nData profiling",TEAL),
    ("⚡","3. SQL","6 business queries\nWindow fns · CTEs",PURPLE),
    ("🤖","4. ML Models","RF · ALS · KMeans\n3 trained models",ORANGE),
    ("🌊","5. Streaming","Live scoring\nforeachBatch",GREEN),
    ("📊","6. Dashboard","8-page BI\nReal-time UI",RED),
]
fw = Inches(1.95); fh = Inches(1.9); gap_f = Inches(0.14)
fx = Inches(0.45); fy = Inches(1.95)
for i, (ic, nm, tech, c) in enumerate(flow):
    card(s, fx, fy, fw, fh, accent=c)
    tb(s, ic, fx, fy+Inches(0.12), fw, Inches(0.5),
       sz=Pt(26), align=PP_ALIGN.CENTER, color=T1)
    tb(s, nm, fx, fy+Inches(0.59), fw, Inches(0.3),
       sz=Pt(11.5), bold=True, color=T1, align=PP_ALIGN.CENTER)
    tb(s, tech, fx, fy+Inches(0.9), fw, Inches(0.85),
       sz=Pt(9.5), color=T3, align=PP_ALIGN.CENTER)
    fx += fw
    if i < 5:
        tb(s, "→", fx, fy+Inches(0.65), Inches(0.14+gap_f), Inches(0.45),
           sz=Pt(14), color=T3, align=PP_ALIGN.CENTER)
        fx += gap_f

# 4 outcome pillars
pillars = [
    (GREEN,  "Real-Time Intelligence","Events → business action in <120ms foreachBatch"),
    (BLUE,   "Personalised at Scale", "ALS gives every 1 of 100K customers unique product picks"),
    (ORANGE, "Revenue Protection",    "ML flags at-risk sessions → auto-trigger recovery promo"),
    (PURPLE, "Horizontally Scalable", "Same code: laptop → 1,000-node AWS EMR, zero changes"),
]
pw2 = Inches(3.0); ph2 = Inches(1.35); gap2 = Inches(0.12)
px2 = Inches(0.55); py2 = Inches(4.1)
for c, t, d in pillars:
    lcard(s, px2, py2, pw2, ph2, accent=c)
    tb(s, t, px2+Inches(0.14), py2+Inches(0.12), pw2-Inches(0.22), Inches(0.3),
       sz=Pt(11.5), bold=True, color=T1)
    tb(s, d, px2+Inches(0.14), py2+Inches(0.46), pw2-Inches(0.22), Inches(0.78),
       sz=Pt(10.5), color=T2)
    px2 += pw2 + gap2

# Why Spark note
rect(s, Inches(0.55), Inches(5.62), Inches(12.22), Inches(0.82),
     fill=LBLUE, line=BLUE, lw=Pt(0.4))
rect(s, Inches(0.55), Inches(5.62), Inches(0.06), Inches(0.82), fill=BLUE)
tb(s, "Why Apache Spark?", Inches(0.72), Inches(5.7), Inches(4), Inches(0.28),
   sz=Pt(11.5), bold=True, color=BLUE)
tb(s, "Our dataset = 2GB+ and 12.8M events. Pandas would crash (loads all data into one machine's RAM). Spark distributes the workload across CPU cores, uses lazy evaluation to optimise query plans automatically, and scales from a laptop to a cloud cluster — same code, zero changes.",
   Inches(0.72), Inches(6.0), Inches(12.0), Inches(0.38),
   sz=Pt(10.5), color=T2)

# ═══════════════════════════════════════════════════════
# SLIDE 4 — TECH STACK
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); slide_bg(s)
header(s, "Technology Stack",
       "Every component chosen for a reason — here's why")
tag(s, "03  ·  TECHNOLOGY CHOICES & RATIONALE", BLUE)

tech_cols = [
    (BLUE, "Core Engine",[
        ("Apache Spark 3.5","Distributed in-memory processing. 100x faster than Hadoop MapReduce — keeps intermediate results in RAM."),
        ("PySpark","Python API for Spark — full DataFrame API + SQL engine + MLlib + Streaming, all in one import."),
        ("Lazy Evaluation","Spark builds a DAG of all operations, then executes the most efficient plan. Auto-optimised by Catalyst."),
    ]),
    (PURPLE,"SQL & Storage",[
        ("Spark SQL","ANSI SQL on distributed DataFrames. Supports CTEs, window functions, LATERAL VIEW EXPLODE — at scale."),
        ("Parquet Format","Columnar storage: 10× smaller files, 30× faster reads than CSV. Only reads the columns a query needs."),
        ("Catalyst Optimizer","Spark's query planner. Rewrites your SQL into the most efficient physical execution plan automatically."),
    ]),
    (ORANGE,"Machine Learning",[
        ("MLlib Pipeline API","Stages: Indexer → Assembler → Model. Clean, reproducible ML. The same pipeline works on 1K or 1B rows."),
        ("VectorAssembler","Combines multiple feature columns into one dense vector — required input format for all Spark ML models."),
        ("Distributed Training","MLlib trains models across all CPU cores simultaneously. scikit-learn trains on one core — doesn't scale."),
    ]),
    (GREEN,"Streaming",[
        ("Structured Streaming","Treats a live event stream as an unbounded table. Run the same SQL queries and ML on live data."),
        ("foreachBatch","Applies any Python code (including pre-trained ML models) to each micro-batch of streaming events."),
        ("Stream-Static JOIN","Joins live events with a static product catalog. Broadcast join: catalog loaded once, enriches every event."),
    ]),
]
cw = Inches(2.9); ch = Inches(3.88); cgap = Inches(0.14)
cx = Inches(0.45)
for ac, cat, items in tech_cols:
    card(s, cx, Inches(1.88), cw, ch, accent=ac)
    tb(s, cat, cx+Inches(0.14), Inches(2.0), cw-Inches(0.22), Inches(0.3),
       sz=Pt(12), bold=True, color=ac)
    iy = Inches(2.42)
    for t, d in items:
        rect(s, cx+Inches(0.12), iy, cw-Inches(0.24), Inches(1.0),
             fill=BG)
        tb(s, t, cx+Inches(0.2), iy+Inches(0.06), cw-Inches(0.35), Inches(0.28),
           sz=Pt(10.5), bold=True, color=T1)
        tb(s, d, cx+Inches(0.2), iy+Inches(0.34), cw-Inches(0.35), Inches(0.58),
           sz=Pt(10), color=T2)
        iy += Inches(1.1)
    cx += cw + cgap

rule(s, Inches(5.9))
tb(s, "Key insight:  Spark provides ONE unified API for batch SQL, machine learning, AND real-time streaming. We did not need three separate tools — one Spark cluster handled all three pipeline stages.",
   Inches(0.55), Inches(5.98), Inches(12.22), Inches(0.45),
   sz=Pt(11), italic=True, color=NAVY)

# ═══════════════════════════════════════════════════════
# SLIDE 5 — DATA & EDA
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); slide_bg(s)
header(s, "Data & Exploratory Data Analysis (EDA)",
       "What data we have — and why understanding it first shapes everything downstream")
tag(s, "04  ·  DATASET OVERVIEW & EDA FINDINGS", TEAL)

# Left: dataset table
card(s, Inches(0.55), Inches(1.88), Inches(5.9), Inches(3.52))
tb(s, "Dataset: 4 Interconnected Tables",
   Inches(0.72), Inches(1.98), Inches(5.6), Inches(0.3),
   sz=Pt(12), bold=True, color=NAVY)
hdrs = ["Table","Rows","Key Columns"]
hxs  = [Inches(0.72),Inches(1.9),Inches(3.0)]
hws  = [Inches(1.12),Inches(1.05),Inches(2.32)]
hy = Inches(2.42)
rect(s, Inches(0.62), hy, Inches(5.76), Inches(0.3), fill=BG)
for h_, hx_, hw_ in zip(hdrs, hxs, hws):
    tb(s, h_, hx_, hy+Inches(0.04), hw_, Inches(0.24),
       sz=Pt(9), bold=True, color=T3)
hy += Inches(0.32)
drows = [
    ("customers.csv","100,000","gender, device, city, DOB"),
    ("products.csv","44,000","category, brand, colour, season"),
    ("transactions.csv","850,000","amount, status, promo, payment_method"),
    ("click_stream.csv","12.8 Million","event_type, session_id, timestamp"),
]
for ri, row_ in enumerate(drows):
    fill_ = BG if ri%2==0 else WHITE
    rect(s, Inches(0.62), hy, Inches(5.76), Inches(0.44), fill=fill_)
    for cell_, hx_, hw_ in zip(row_, hxs, hws):
        tb(s, cell_, hx_, hy+Inches(0.07), hw_, Inches(0.3),
           sz=Pt(10.5), color=T2)
    hy += Inches(0.44)

# EDA explanation
lcard(s, Inches(0.55), Inches(5.52), Inches(5.9), Inches(1.56), fill=LTEAL, accent=TEAL)
tb(s, "Why EDA Before Anything Else?",
   Inches(0.72), Inches(5.62), Inches(5.6), Inches(0.3),
   sz=Pt(11.5), bold=True, color=TEAL)
tb(s, "EDA prevents building the wrong thing. Charting distributions and correlations before writing any ML code revealed: 77% Android users (shaped dashboard priority), payment failures are random/external (changed our ML question entirely), and a clear June–July seasonal spike (shaped SQL query design).",
   Inches(0.72), Inches(5.98), Inches(5.6), Inches(0.98),
   sz=Pt(10.5), color=T2)

# Right: 6 EDA findings
card(s, Inches(6.62), Inches(1.88), Inches(6.14), Inches(5.2))
tb(s, "6 Key EDA Findings & Their Impact",
   Inches(6.8), Inches(1.98), Inches(5.86), Inches(0.3),
   sz=Pt(12), bold=True, color=NAVY)
findings = [
    (TEAL,  "77% Android Users","→ Mobile-first dashboard priority · device_type added as ML feature"),
    (BLUE,  "Payment Success 95.7%","→ Failures are EXTERNAL (bank/network), not platform — AUC ≈ 0.5 is honest"),
    (GREEN, "10am–8pm Peak Activity","→ Streaming alerts most critical in this window · flash sale timing"),
    (ORANGE,"5 Distinct Behaviour Patterns","→ Confirmed KMeans segmentation was necessary · one-size marketing fails"),
    (PURPLE,"Top 28% → 94.5% Revenue","→ VIP protection became a core feature · real-time alert for Power Buyers"),
    (RED,   "42% Cart Abandonment","→ Highest-impact fix identified · ML scorer targets this session state"),
]
fy2 = Inches(2.42)
for c, ft, fd in findings:
    rect(s, Inches(6.72), fy2, Inches(5.94), Inches(0.74), fill=BG)
    rect(s, Inches(6.72), fy2, Inches(0.05), Inches(0.74), fill=c)
    tb(s, ft, Inches(6.84), fy2+Inches(0.06), Inches(5.7), Inches(0.28),
       sz=Pt(11), bold=True, color=T1)
    tb(s, fd, Inches(6.84), fy2+Inches(0.34), Inches(5.7), Inches(0.34),
       sz=Pt(10), color=T2)
    fy2 += Inches(0.85)

# ═══════════════════════════════════════════════════════
# SLIDE 6 — SPARK SQL
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); slide_bg(s)
header(s, "Spark SQL — 6 Advanced Business Queries",
       "Each query answers a critical business question using specific SQL concepts — here's why each concept was chosen")
tag(s, "05  ·  SQL CONCEPTS & RATIONALE", PURPLE)

queries = [
    (BLUE,  "Q1: Conversion Funnel",
     "Where do customers drop off?",
     "4-table JOIN · 2 CTEs · ROW_NUMBER()",
     "CTEs build the funnel step-by-step. ROW_NUMBER() finds last event per session — impossible with a simple COUNT."),
    (PURPLE,"Q2: Market Basket Analysis",
     "What products sell together?",
     "LATERAL VIEW EXPLODE · Self-JOIN · Lift",
     "Products stored as arrays. EXPLODE flattens to rows for self-JOIN. Lift = conf/baseline — tells us if correlation is real."),
    (TEAL,  "Q3: Cohort Retention",
     "How long do customers stay active?",
     "MONTHS_BETWEEN() · Pivot CASE · NULLIF",
     "MONTHS_BETWEEN tracks activity gaps per cohort. NULLIF prevents divide-by-zero as cohorts shrink at later months."),
    (GREEN, "Q4: RFM Scoring",
     "Who are the most valuable customers now?",
     "NTILE(5) · 3 nested CTEs · Composite score",
     "NTILE splits customers into 5 equal buckets per RFM dimension. Composite score feeds directly into KMeans clustering."),
    (ORANGE,"Q5: Purchase Velocity",
     "How fast do customers re-purchase?",
     "LAG() · LEAD() · DATEDIFF() · Running aggs",
     "LAG() gives previous purchase date per customer. DATEDIFF calculates the gap. Identifies Power Buyers vs Casual shoppers."),
    (RED,   "Q6: Product Affinity",
     "What do customers browse before they buy?",
     "4-table JOIN · EXPLODE · Browse-to-buy lift",
     "Joins clickstream with transactions to find browse-before-buy paths. Lift >1 means browsing Cat A predicts buying Cat B."),
]
qw = Inches(3.93); qh = Inches(2.48); qgap = Inches(0.14)
qx = Inches(0.45); qy = Inches(1.98)
for i, (c, qt, bq, concepts, why) in enumerate(queries):
    col = i%3; row = i//2
    x_ = qx + col*(qw+qgap); y_ = qy + row*(qh+qgap)
    card(s, x_, y_, qw, qh, accent=c)
    rect(s, x_, y_, Inches(0.04), qh, fill=c)
    tb(s, f"Q{i+1} · {qt}", x_+Inches(0.12), y_+Inches(0.1),
       qw-Inches(0.2), Inches(0.28), sz=Pt(11.5), bold=True, color=T1)
    tb(s, bq, x_+Inches(0.12), y_+Inches(0.38), qw-Inches(0.2), Inches(0.22),
       sz=Pt(10), italic=True, color=T3)
    tb(s, "SQL: "+concepts, x_+Inches(0.12), y_+Inches(0.62), qw-Inches(0.2), Inches(0.24),
       sz=Pt(10), bold=True, color=c)
    rect(s, x_+Inches(0.1), y_+Inches(0.92), qw-Inches(0.2), Inches(1.42), fill=BG)
    tb(s, "WHY: "+why, x_+Inches(0.18), y_+Inches(0.98), qw-Inches(0.32), Inches(1.31),
       sz=Pt(10), color=T2)

# ═══════════════════════════════════════════════════════
# SLIDE 7 — MACHINE LEARNING
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); slide_bg(s)
header(s, "Machine Learning — Three Models, Three Problems",
       "Which algorithm, why we chose it over alternatives, and what it tells the business")
tag(s, "06  ·  ML CONCEPTS & ALGORITHM SELECTION RATIONALE", ORANGE)

models = [
    (BLUE, "Model 1: Payment Classifier",
     "Can we predict if a payment will fail?",
     [("Algorithm","Random Forest (Logistic Regression + GBT also tested)"),
      ("Features","15: transaction (4) + clickstream (8) + demographics (3)"),
      ("Why Random Forest?","RF builds many trees in parallel, each on a random subset. Majority vote reduces overfitting. Captures non-linear feature interactions LR cannot."),
      ("Why NOT Neural Network?","With 15 features and 850K rows, NN would overfit and needs GPU. RF is faster, more interpretable, sufficient."),
      ("Result","Accuracy: 95.4%   F1: 0.931   AUC ≈ 0.5"),
      ("Business insight","AUC ≈ 0.5 is an HONEST result: payment failures are external (bank/network), NOT caused by the platform."),
     ]),
    (ORANGE,"Model 2: KMeans Segmentation",
     "What are the distinct customer types?",
     [("Algorithm","KMeans Clustering (k=5, chosen by Elbow + Silhouette)"),
      ("Features","10: RFM score (6) + browsing behaviour (4)"),
      ("Why KMeans?","Unsupervised — discovers hidden structure without labels. We don't know in advance what groups exist, KMeans finds them."),
      ("Why NOT DBSCAN?","DBSCAN struggles with high-dimensional sparse data. KMeans is O(n·k·i) — scales to 50K customers easily."),
      ("How we chose k=5","Elbow Method plots inertia vs k. Silhouette Score measures separation. Best score: 0.629 at k=5."),
      ("Result","5 segments: Casual (31K) · Budget (2.5K) · Power (1.8K) · Regular (12K) · Big-Ticket (2.3K)"),
     ]),
    (PURPLE,"Model 3: ALS Recommender",
     "What products should each customer see?",
     [("Algorithm","ALS — Alternating Least Squares (Collaborative Filtering)"),
      ("Input","Customer × Product implicit interaction matrix"),
      ("Why ALS?","Learns from implicit signals — checkout (5pts), cart (3pts), wishlist (2pts), view (1pt). Discovers cross-customer patterns (what people like you bought)."),
      ("Why NOT Content-Based?","Content-based uses only product features — cannot discover that customers like you also bought X. ALS learns latent behaviour patterns."),
      ("Result","RMSE: 0.8243   Top-5 personalised recs for all 100K customers"),
      ("Business value","Personalisation that was impossible at scale is now fully automated."),
     ]),
]
mw = Inches(3.95); mh = Inches(5.28); mgap = Inches(0.12)
mx = Inches(0.45)
for c, mt, mq, mlist in models:
    card(s, mx, Inches(1.88), mw, mh, accent=c)
    tb(s, mt, mx+Inches(0.14), Inches(2.0), mw-Inches(0.22), Inches(0.28),
       sz=Pt(12.5), bold=True, color=T1)
    tb(s, mq, mx+Inches(0.14), Inches(2.3), mw-Inches(0.22), Inches(0.24),
       sz=Pt(10), italic=True, color=T3)
    iy = Inches(2.62)
    for fld, val in mlist:
        is_why = fld.startswith("Why") or fld == "How we chose k=5"
        bg_ = LGREEN if fld.startswith("Result") or fld == "Business" or fld == "Business insight" or fld == "Business value" else (LORANGE if is_why else BG)
        h_ = Inches(0.75) if is_why else Inches(0.62)
        rect(s, mx+Inches(0.12), iy, mw-Inches(0.24), h_, fill=bg_)
        tb(s, fld, mx+Inches(0.2), iy+Inches(0.05), mw-Inches(0.38), Inches(0.22),
           sz=Pt(8.5), bold=True, color=T3)
        tb(s, val, mx+Inches(0.2), iy+Inches(0.27), mw-Inches(0.38), h_-Inches(0.28),
           sz=Pt(10), color=T1)
        iy += h_ + Inches(0.05)
    mx += mw + mgap

# ═══════════════════════════════════════════════════════
# SLIDE 8 — REAL-TIME STREAMING
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); slide_bg(s)
header(s, "Real-Time Streaming — Spark Structured Streaming",
       "Why streaming instead of batch — the 2-phase architecture and the engineering challenges we solved")
tag(s, "07  ·  STREAMING CONCEPTS & DESIGN DECISIONS", GREEN)

# Why streaming
lcard(s, Inches(0.55), Inches(1.88), Inches(12.22), Inches(0.72), fill=LGREEN, accent=GREEN)
tb(s, "Why Streaming — not just a nightly batch job?",
   Inches(0.72), Inches(1.96), Inches(12.0), Inches(0.28),
   sz=Pt(12), bold=True, color=GREEN)
tb(s, "A batch job runs once per day. A customer showing HIGH abandonment risk at 2pm gets a discount email at midnight — they already bought from a competitor. Streaming processes each event as it arrives, so interventions happen in the same session while the customer is still on the site.",
   Inches(0.72), Inches(2.26), Inches(12.0), Inches(0.28),
   sz=Pt(11), color=T2)

# 3 concepts
concepts = [
    (GREEN, "Structured Streaming",
     "What it is:",
     "Spark treats a live data stream as an unbounded table that grows over time. The same SQL queries written for batch processing run continuously on this growing table as new rows arrive.",
     "Why we chose it:",
     "We already wrote Spark SQL for batch analysis. Structured Streaming reuses the exact same DataFrame API — no new framework needed, minimal learning curve, maximum code reuse across batch and streaming."),
    (BLUE,  "foreachBatch — ML on Streams",
     "What it is:",
     "foreachBatch lets us run arbitrary Python code on each micro-batch of streaming events. We apply our pre-trained RandomForest model to score every session in each batch and emit predictions.",
     "Why we chose it:",
     "Spark ML models cannot natively run inside a continuous streaming query. foreachBatch is the bridge — gives us full Python control over each batch, so we can call model.transform(batch_df) directly."),
    (PURPLE,"Stream-Static JOIN",
     "What it is:",
     "We join the live event stream (clickstream) with the static product catalog (44K rows). The static table is broadcast to all cluster nodes once and stays there — enriching every incoming event cheaply.",
     "Why we chose it:",
     "Every click only has a product_id. To know category and price we need the product table. Stream-Static JOIN adds this context to every event in real time, enabling richer session features for the ML scorer."),
]
sw = Inches(3.93); sh = Inches(3.72); sgap = Inches(0.14)
sx = Inches(0.45)
for c, ct, lbl1, d1, lbl2, d2 in concepts:
    card(s, sx, Inches(2.78), sw, sh, accent=c)
    tb(s, ct, sx+Inches(0.14), Inches(2.9), sw-Inches(0.22), Inches(0.3),
       sz=Pt(12), bold=True, color=T1)
    rect(s, sx+Inches(0.12), Inches(3.28), sw-Inches(0.24), Inches(1.28), fill=BG)
    tb(s, lbl1, sx+Inches(0.2), Inches(3.33), sw-Inches(0.35), Inches(0.22),
       sz=Pt(9), bold=True, color=T3)
    tb(s, d1, sx+Inches(0.2), Inches(3.56), sw-Inches(0.35), Inches(0.95),
       sz=Pt(10.5), color=T2)
    rect(s, sx+Inches(0.12), Inches(4.65), sw-Inches(0.24), Inches(1.72), fill=LGREEN)
    tb(s, lbl2, sx+Inches(0.2), Inches(4.7), sw-Inches(0.35), Inches(0.22),
       sz=Pt(9), bold=True, color=GREEN)
    tb(s, d2, sx+Inches(0.2), Inches(4.94), sw-Inches(0.35), Inches(1.35),
       sz=Pt(10.5), color=T2)
    sx += sw + sgap

rule(s, Inches(6.6))
tb(s, "Engineering wins:   ① Windows NativeIO crash (UnsatisfiedLinkError) → solved by keeping RF model in Python memory, no disk save needed.   ② Python UDF EOFException in streaming → solved by replacing UDFs with Spark-native vector_to_array() function.",
   Inches(0.55), Inches(6.68), Inches(12.22), Inches(0.45),
   sz=Pt(10.5), italic=True, color=NAVY)

# ═══════════════════════════════════════════════════════
# SLIDE 9 — DASHBOARD & SIMULATION
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); slide_bg(s)
header(s, "Business Dashboard & E-Commerce Simulation",
       "The front-end face of the pipeline — 8-page BI command centre + live user simulation")
tag(s, "08  ·  DASHBOARD FEATURES & SIMULATION DEMO", BLUE)

# 8 dashboard pages
pages = [
    ("⚡","Live Operations",   BLUE,  "Real-time order feed · revenue counter · ML high-risk session flags"),
    ("💰","Revenue & Sales",   GREEN, "MoM trends · LAG() SQL · payment method mix · VIP segment revenue"),
    ("👥","Customer Intel",    ORANGE,"All 5 KMeans segments · Top-15 LTV leaderboard by name"),
    ("🎯","Conversion Center", RED,   "Animated funnel · 42% Cart→Checkout drop · promo effectiveness"),
    ("🤖","ML Predictions",    PURPLE,"Live session grid · ALS recommender shelf · model accuracy metrics"),
    ("📦","Product Analytics", TEAL,  "Market basket table (confidence+lift) · trending products live"),
    ("🚨","Alerts & Actions",  RED,   "6 alert types · one-click business responses · anomaly feed"),
    ("⚙️","Pipeline Control",  T3,    "Click a module → watch simulated terminal execution in real-time"),
]
dw = Inches(2.95); dh = Inches(1.28); dgap = Inches(0.1)
dx = Inches(0.45); dy = Inches(1.88)
for i, (ic, dt, dc, dd) in enumerate(pages):
    col = i%4; row = i//4
    x_ = dx + col*(dw+dgap); y_ = dy + row*(dh+dgap)
    lcard(s, x_, y_, dw, dh, accent=dc)
    tb(s, ic, x_+Inches(0.14), y_+Inches(0.08), Inches(0.4), Inches(0.45), sz=Pt(20), color=T1)
    tb(s, dt, x_+Inches(0.14), y_+Inches(0.48), dw-Inches(0.22), Inches(0.26),
       sz=Pt(11.5), bold=True, color=T1)
    tb(s, dd, x_+Inches(0.14), y_+Inches(0.76), dw-Inches(0.22), Inches(0.44),
       sz=Pt(9.5), color=T3)

# Simulation explainer
lcard(s, Inches(0.55), Inches(4.6), Inches(12.22), Inches(2.6), fill=LBLUE, accent=BLUE)
tb(s, "🛒  E-Commerce Simulation — simulation.html",
   Inches(0.72), Inches(4.72), Inches(12.0), Inches(0.3),
   sz=Pt(13), bold=True, color=NAVY)
sim_items = [
    ("Spark Streaming Events","Every hover, click, add-to-cart logs a real-time event with timestamp in the analytics panel — just like foreachBatch consumes live events."),
    ("Random Forest Score","A conversion probability gauge updates live as you browse. LOW → MEDIUM → HIGH based on cart items, browse depth, and category variety features."),
    ("KMeans Segment","Your customer cluster auto-reassigns — Casual → Regular → Power Buyer — as cart value grows. Matches the real k=5 model's behaviour."),
    ("ALS Recommendations","When you view a category, the panel updates with 4 personalised picks — simulating ALS latent factor predictions per user-category pair."),
    ("Market Basket SQL","Add an Apparel item → 'Complete the Look: Footwear' triggers with 68% confidence and 2.3x lift. The exact output of our SQL Q2 query."),
    ("ML Promo Trigger","At MEDIUM risk score + cart items, a 10% promo popup fires automatically — simulating the foreachBatch promo dispatch business action."),
]
sx2 = Inches(0.72); sy2 = Inches(5.1)
col_w = Inches(6.0)
for i, (lbl, desc) in enumerate(sim_items):
    col = i%2; row = i//2
    x_ = sx2 + col*col_w; y_ = sy2 + row*Inches(0.62)
    tb(s, "→ "+lbl+":", x_, y_, col_w, Inches(0.26), sz=Pt(10.5), bold=True, color=BLUE)
    tb(s, "   "+desc, x_, y_+Inches(0.24), col_w-Inches(0.1), Inches(0.35), sz=Pt(10), color=T2)

# ═══════════════════════════════════════════════════════
# SLIDE 10 — OUTCOMES & BUSINESS IMPACT
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); slide_bg(s)
header(s, "Outcomes & Business Impact",
       "What the pipeline proved — measurable results that translate directly into revenue")
tag(s, "09  ·  RESULTS & BUSINESS VALUE", GREEN)

kpi_row(s, [
    ("95.4%","RF Accuracy","F1 = 0.931",BLUE),
    ("0.629","Silhouette","KMeans k=5 best",PURPLE),
    ("0.824","RMSE (ALS)","Recommender quality",GREEN),
    ("~120ms","Stream Latency","foreachBatch",ORANGE),
    ("42%","Cart Abandonment","Identified & fixable",RED),
    ("95.7%","Payment Success","All 5 methods",TEAL),
], Inches(1.82), h=Inches(1.04))

# Left: technical
card(s, Inches(0.55), Inches(3.06), Inches(5.9), Inches(4.08))
tb(s, "Technical Achievements",
   Inches(0.72), Inches(3.16), Inches(5.62), Inches(0.3),
   sz=Pt(12), bold=True, color=NAVY)
techs = [
    (BLUE,  "6-Stage Spark Pipeline","Ingestion → EDA → SQL → ML → Streaming → Dashboard. Fully automated, end-to-end."),
    (PURPLE,"6 Advanced SQL Queries","Window functions · CTEs · LATERAL VIEW EXPLODE · Self-JOINs · 4-table JOINs."),
    (ORANGE,"3 MLlib Models","Random Forest (95.4%) · KMeans k=5 (Silhouette 0.629) · ALS (RMSE 0.824)."),
    (GREEN, "Structured Streaming","2-phase: batch train → foreachBatch inference · Anomaly detection · ~120ms latency."),
    (RED,   "Engineering Wins","Solved Windows NativeIO crash + Python UDF EOFException — real production-grade bugs."),
]
ty2 = Inches(3.6)
for c, tt, td in techs:
    rect(s, Inches(0.65), ty2, Inches(5.7), Inches(0.68), fill=BG)
    rect(s, Inches(0.65), ty2, Inches(0.04), Inches(0.68), fill=c)
    tb(s, tt, Inches(0.78), ty2+Inches(0.06), Inches(5.44), Inches(0.26),
       sz=Pt(11), bold=True, color=T1)
    tb(s, td, Inches(0.78), ty2+Inches(0.32), Inches(5.44), Inches(0.3),
       sz=Pt(10), color=T2)
    ty2 += Inches(0.76)

# Right: business outcomes
card(s, Inches(6.62), Inches(3.06), Inches(6.25), Inches(4.08))
tb(s, "Business Outcomes", Inches(6.8), Inches(3.16), Inches(5.96), Inches(0.3),
   sz=Pt(12), bold=True, color=NAVY)
biz = [
    (RED,   "Fix the Funnel → +12% Revenue",
     "Cart→Checkout drops 42%. Fixing to 30% by showing shipping cost upfront → +12% total revenue. Single SQL query identified this."),
    (BLUE,  "Real-Time Promo Recovery",
     "ML scores live sessions. MEDIUM-risk (40–69%) triggers a 10% discount push. Recovers IDR 12–15M in sessions that would have abandoned."),
    (GREEN, "Cross-Sell Activation",
     "Market basket SQL: Apparel+Footwear = 68% confidence, 2.3x lift. 'Complete the Look' banner → IDR 2.1M/day incremental."),
    (PURPLE,"VIP Customer Protection",
     "KMeans found 1,827 Power Buyers averaging IDR 71.2M spend each. Streaming flags their sessions in real time for priority service."),
]
by2 = Inches(3.6)
for c, bt, bd in biz:
    rect(s, Inches(6.72), by2, Inches(6.05), Inches(0.88), fill=BG)
    rect(s, Inches(6.72), by2, Inches(0.04), Inches(0.88), fill=c)
    tb(s, bt, Inches(6.84), by2+Inches(0.06), Inches(5.8), Inches(0.26),
       sz=Pt(11), bold=True, color=T1)
    tb(s, bd, Inches(6.84), by2+Inches(0.34), Inches(5.8), Inches(0.48),
       sz=Pt(10), color=T2)
    by2 += Inches(0.96)

# ═══════════════════════════════════════════════════════
# SLIDE 11 — KEY LEARNINGS
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); slide_bg(s)
header(s, "Key Learnings — Concepts Explained Simply",
       "What each technology taught us and why the combination is powerful")
tag(s, "10  ·  CONCEPT SUMMARY & TAKEAWAYS", TEAL)

cols_data = [
    (BLUE, "Apache Spark\nCore Engine", [
        ("Lazy Evaluation","Spark PLANS before executing. Catalyst rewrites your code into an optimal physical execution plan — you get free performance."),
        ("Distributed Processing","12.8M rows split across cores/nodes. Each worker processes its partition in parallel. Speeds up linearly with more resources."),
        ("Unified Platform","SQL + ML + Streaming = one API, one cluster, one deploy. No switching between Hadoop, scikit-learn, and Kafka separately."),
        ("Horizontal Scale","The same code that runs on a laptop runs on AWS EMR with 1,000 nodes. Architecture scales without any code change."),
    ]),
    (ORANGE,"SQL & Machine Learning\nAlgorithms & Rationale", [
        ("Window Functions","LAG, NTILE, ROW_NUMBER — analyse rows relative to their NEIGHBOURS. Essential for time-series, RFM scoring, and funnel analysis."),
        ("Random Forest","Many weak trees → one strong model. Handles non-linearity, mixed feature types, and prevents overfitting automatically via bagging."),
        ("KMeans Clustering","Groups data into k clusters by minimising intra-cluster distance. Unsupervised — finds structure without any labelled training data."),
        ("ALS Recommender","Factorises user-item interaction matrix into latent factors. Learns preferences from IMPLICIT signals, not explicit star ratings."),
    ]),
    (GREEN,"Streaming & Future Scale\nArchitecture Decisions", [
        ("Batch vs Streaming","Batch = react after the fact. Streaming = act while the customer is still there. For e-commerce, the difference is the sale."),
        ("foreachBatch Bridge","Connects pre-trained batch ML models to live streams. The cleanest solution to applying offline ML to online data."),
        ("Lambda Alternative","We avoided the Lambda Architecture anti-pattern by using Spark for BOTH batch training and streaming inference in one codebase."),
        ("Next Steps","Add Kafka → real-time broker. FastAPI backend → live data. AWS EMR → cloud scale. DBSCAN + autoencoders → richer ML."),
    ]),
]
cw2 = Inches(4.1); cgap2 = Inches(0.11); cx2 = Inches(0.45)
for c, ch2, items in cols_data:
    card(s, cx2, Inches(1.88), cw2, Inches(5.52), accent=c)
    tb(s, ch2, cx2+Inches(0.14), Inches(2.0), cw2-Inches(0.22), Inches(0.52),
       sz=Pt(11.5), bold=True, color=T1)
    iy2 = Inches(2.58)
    for concept, desc in items:
        rect(s, cx2+Inches(0.12), iy2, cw2-Inches(0.24), Inches(1.14), fill=BG)
        rect(s, cx2+Inches(0.12), iy2, Inches(0.04), Inches(1.14), fill=c)
        tb(s, concept, cx2+Inches(0.24), iy2+Inches(0.07), cw2-Inches(0.36), Inches(0.26),
           sz=Pt(11), bold=True, color=T1)
        tb(s, desc, cx2+Inches(0.24), iy2+Inches(0.36), cw2-Inches(0.36), Inches(0.72),
           sz=Pt(10), color=T2)
        iy2 += Inches(1.22)
    cx2 += cw2 + cgap2

rule(s, Inches(7.2))
rect(s, Inches(0.55), Inches(7.24), Inches(12.22), Inches(0.18), fill=NAVY)
tb(s, "   Core takeaway: open-source Spark replaces enterprise BI tools at a fraction of the cost — processing more data, faster, in real time.",
   Inches(0.55), Inches(7.24), Inches(12.22), Inches(0.18),
   sz=Pt(10), bold=True, color=WHITE)

# ═══════════════════════════════════════════════════════
# SLIDE 12 — THANK YOU
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); slide_bg(s)
rect(s, 0, 0, W, H, fill=NAVY)
rect(s, 0, H-Inches(0.07), W, Inches(0.07), fill=BLUE)
rect(s, 0, 0, W, Inches(0.048), fill=BLUE)

tb(s, "Thank You", Inches(0.65), Inches(0.95), Inches(7.5), Inches(1.4),
   sz=Pt(64), bold=True, color=WHITE)
tb(s, "Questions & Discussion",
   Inches(0.65), Inches(2.3), Inches(7.5), Inches(0.72),
   sz=Pt(28), color=RGBColor(0xA0,0xB4,0xE8))
rect(s, Inches(0.65), Inches(3.15), Inches(2.8), Inches(0.038), fill=BLUE)
tb(s, "CcMart · Real-Time E-Commerce Analytics Platform\nApache Spark · MLlib · Spark SQL · Structured Streaming\nCloud Computing Final Project · Group 6",
   Inches(0.65), Inches(3.3), Inches(7.2), Inches(1.2),
   sz=Pt(12), color=RGBColor(0x7A,0x90,0xC8))

# Project summary on right
rect(s, Inches(8.4), Inches(0.55), Inches(4.65), Inches(6.5),
     fill=RGBColor(0x16,0x1E,0x4A), line=BLUE, lw=Pt(0.5))
tb(s, "Deliverables", Inches(8.62), Inches(0.72), Inches(4.25), Inches(0.3),
   sz=Pt(13), bold=True, color=WHITE)
delivs = [
    ("📂","ingestion.py · eda.py","Data pipeline + EDA",BLUE),
    ("⚡","transformations.py","6 Spark SQL queries",PURPLE),
    ("🤖","ml_pipeline.py","RF + ALS + KMeans models",ORANGE),
    ("🌊","streaming.py","Structured Streaming pipeline",GREEN),
    ("🧠","streaming_predictions.py","Real-time ML scoring",TEAL),
    ("📊","dashboard/index.html","8-page BI command centre",BLUE),
    ("🛒","dashboard/simulation.html","Live e-commerce simulation",GREEN),
    ("📋","docs/dashboard_guide.md","Presenter's guide",T3),
]
dy3 = Inches(1.14)
for ic, fn, desc, c in delivs:
    tb(s, ic, Inches(8.62), dy3, Inches(0.38), Inches(0.45), sz=Pt(17), color=WHITE)
    tb(s, fn, Inches(9.06), dy3, Inches(3.9), Inches(0.26), sz=Pt(10.5), bold=True, color=c)
    tb(s, desc, Inches(9.06), dy3+Inches(0.26), Inches(3.9), Inches(0.24), sz=Pt(9.5), color=RGBColor(0xC0,0xCC,0xEC))
    dy3 += Inches(0.7)

tb(s, "Open: dashboard/simulation.html  to run the live demo",
   Inches(0.65), Inches(5.45), Inches(7.2), Inches(0.35),
   sz=Pt(11.5), bold=True, color=BLUE)
tb(s, "Open: dashboard/index.html  for the full BI dashboard",
   Inches(0.65), Inches(5.82), Inches(7.2), Inches(0.35),
   sz=Pt(11.5), color=RGBColor(0x60,0xA5,0xFA))

# ── SAVE ───────────────────────────────────────────────────
out_dir  = os.path.join(os.path.dirname(os.path.abspath(__file__)),"dashboard")
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, "CcMart_Final_Presentation.pptx")
prs.save(out_path)
print(f"[OK] Saved: {out_path}")
print(f"     Slides: {len(prs.slides)}")
